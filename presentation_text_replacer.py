# -*- coding: utf-8 -*-
from pptx import Presentation
import pprint as pp
import ConfigParser    
import os
import sys

from arguments import *


class pptRep():
    def __init__(self, inFile, outfile, repFile):
        inFile = os.path.abspath(inFile)                
        if not os.path.isfile(inFile):
            exit('Input file does not exists: %s' % inFile)        
        self.prs = Presentation(inFile)        
        self.loadReplaceFile(repFile)
        self.doReplaces()
        self.savePres(outfile)
    
    def loadReplaceFile(self, repFile):
        repFile = os.path.abspath(repFile)                
        if not os.path.isfile(repFile):        
            exit('File with string to replace doesnot exists: %s' % repFile)                    
        self.cp = ConfigParser.ConfigParser()                
        self.cp.read(repFile)                    

    # retrieves a configuration value from config. file or the same text if it is not found
    def getReplace(self, txtToReplace):
        # the section name does not matter. use the first
        if len(self.cp.sections()):
            section = self.cp.sections()[0]

        searchTxt = txtToReplace.replace('{','').replace('}','').strip()
        if self.cp.has_option(section, searchTxt):
            rep = self.cp.get(section, searchTxt)            
            return rep
        
        return txtToReplace
        
    def savePres(self, outputFile):
        outputFile = os.path.abspath(outputFile)                
        self.prs.save(outputFile)
        
    
    def doReplaces(self):        
        '''
        Replaces each ocurrence of text inside presentation file.
        '''
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pp.pprint(run.text)
                        # replaces if fould, leave it if not found!
                        run.text = self.getReplace(run.text)
                        

if __name__ == "__main__":    
    ppt = pptRep(args.inF, args.outF, args.repF)    
    ppt.doReplaces()
    
