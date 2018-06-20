"""
Replace a list of text strings on presentations
"""
import os
import sys
import pprint as pp
import configparser

from pptx import Presentation

class PptRep():
    def __init__(self, input_file, output_file, replacer_file):
        """
        :param input_file: Input presentation file
        :param output_file: Output presentation file.
        :param replacer_file: Txt (ini format) file with strings to be replaced
        """
        input_file = os.path.abspath(input_file)
        if not os.path.isfile(input_file):
            exit('Input file does not exists: %s' % input_file)
        self.input_file = input_file
        self.output_file = output_file
        self.replacer_file = replacer_file

    def process_file(self):
        self.prs = Presentation(self.input_file)
        self._load_replace_file()
        self._replace_inside_file()
        self._save_presentation_file()

    def _load_replace_file(self):
        f = os.path.abspath(self.replacer_file)
        if not os.path.isfile(f):
            exit('File with string to replace does not exists: %s' % self.replacer_file)
        self.cp = configparser.ConfigParser()
        self.cp.read(f)

    # retrieves a configuration value from config. file or the same text if it is not found
    def _get_replace(self, text_to_replace):
        # the section name does not matter. use the first
        if len(self.cp.sections()):
            section = self.cp.sections()[0]

        search_text = text_to_replace.replace('{', '').replace('}', '').strip()
        if self.cp.has_option(section, search_text):
            rep = self.cp.get(section, search_text)
            return rep
        
        return text_to_replace
        
    def _save_presentation_file(self):
        self.prs.save(self.output_file)
        
    
    def _replace_inside_file(self):
        '''
        Replaces each occurrence of text inside presentation file.
        '''
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if '{' in run.text:
                            if '}' in run.text:
                                pp.pprint(run.text)
                                # replaces if fould, leave it if not found!
                                run.text = self._get_replace(run.text)
